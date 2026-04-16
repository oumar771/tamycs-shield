"""
Génère le fichier TAMYCS_Shield_Presentation.pptx
Exécuter : python rapport/make_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Couleurs de la charte ──────────────────────────────────────────────────
BLUE      = RGBColor(0x1A, 0x56, 0xA0)   # bleu principal
BLUE_DARK = RGBColor(0x0D, 0x2B, 0x55)   # bleu foncé (titres slides)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF5, 0xFF)   # fond clair bleuté
GRAY      = RGBColor(0x55, 0x55, 0x55)
GREEN     = RGBColor(0x2E, 0x7D, 0x32)
ORANGE    = RGBColor(0xE6, 0x5C, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # layout vide


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # 1 = MSO_SHAPE_TYPE.RECTANGLE
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_bullet_box(slide, items, x, y, w, h,
                   font_size=16, color=BLUE_DARK, bullet="▸ "):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=None):
    """Barre de titre bleue en haut de chaque slide."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), fill=BLUE_DARK)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.08), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.72), Inches(12), Inches(0.45),
                    font_size=15, color=RGBColor(0xBB, 0xCC, 0xFF))


def accent_line(slide, y=Inches(1.25)):
    """Ligne d'accentuation dorée sous le header."""
    bar = slide.shapes.add_shape(1, 0, y, SLIDE_W, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xFF, 0xC1, 0x07)
    bar.line.fill.background()


def footer(slide, text="TAMYCS Shield — ESAIP La Salle · Bachelor 3 Cybersécurité · 2025-2026"):
    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.4), fill=BLUE_DARK)
    add_textbox(slide, text,
                Inches(0.3), Inches(7.12), Inches(12.5), Inches(0.35),
                font_size=10, color=RGBColor(0xAA, 0xBB, 0xDD),
                align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Couverture
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)

# Fond dégradé simulé (deux rectangles)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BLUE_DARK)
add_rect(slide, 0, Inches(4.5), SLIDE_W, Inches(3), fill=BLUE)

# Titre principal
add_textbox(slide,
            "🛡  TAMYCS Shield",
            Inches(1), Inches(1.2), Inches(11), Inches(1.4),
            font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide,
            "PassVault — Gestionnaire de mots de passe sécurisé",
            Inches(1), Inches(2.6), Inches(11), Inches(0.8),
            font_size=22, color=RGBColor(0xBB, 0xCC, 0xFF),
            align=PP_ALIGN.CENTER, italic=True)

# Bandeau info
add_rect(slide, Inches(2.5), Inches(3.6), Inches(8.3), Inches(0.9),
         fill=RGBColor(0x0A, 0x1F, 0x44))

add_textbox(slide,
            "Conforme ANSSI  ·  AES-256-GCM  ·  JWT  ·  bcrypt",
            Inches(2.5), Inches(3.65), Inches(8.3), Inches(0.75),
            font_size=15, color=RGBColor(0xFF, 0xC1, 0x07),
            align=PP_ALIGN.CENTER, bold=True)

# Auteurs + infos
add_textbox(slide,
            "Oumar Touré  ·  Rodney  ·  Fresnel",
            Inches(1), Inches(5.0), Inches(11), Inches(0.6),
            font_size=18, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

add_textbox(slide,
            "ESAIP La Salle  ·  Bachelor 3 Cybersécurité  ·  2025 / 2026",
            Inches(1), Inches(5.65), Inches(11), Inches(0.5),
            font_size=14, color=RGBColor(0xAA, 0xBB, 0xDD),
            align=PP_ALIGN.CENTER)

add_textbox(slide,
            "Web Programming — Projet de cours",
            Inches(1), Inches(6.2), Inches(11), Inches(0.45),
            font_size=13, color=RGBColor(0x88, 0xAA, 0xCC),
            align=PP_ALIGN.CENTER, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Sommaire
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_BG)
header_bar(slide, "Sommaire", "Plan de présentation")
accent_line(slide)
footer(slide)

items_left = [
    "01 · Introduction & Contexte",
    "02 · Objectifs du projet",
    "03 · Choix technologiques",
    "04 · Architecture générale",
    "05 · Base de données",
]
items_right = [
    "06 · Fonctionnalités clés",
    "07 · Sécurité & ANSSI",
    "08 · Générateur de mots de passe",
    "09 · Gamification",
    "10 · Tests & Conclusion",
]

for i, item in enumerate(items_left):
    y = Inches(1.55) + i * Inches(0.9)
    add_rect(slide, Inches(0.5), y, Inches(5.7), Inches(0.75),
             fill=BLUE_DARK)
    add_textbox(slide, item, Inches(0.6), y + Pt(8),
                Inches(5.5), Inches(0.65),
                font_size=16, color=WHITE, bold=(i == 0))

for i, item in enumerate(items_right):
    y = Inches(1.55) + i * Inches(0.9)
    add_rect(slide, Inches(6.8), y, Inches(5.7), Inches(0.75),
             fill=BLUE)
    add_textbox(slide, item, Inches(6.9), y + Pt(8),
                Inches(5.5), Inches(0.65),
                font_size=16, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Introduction & Contexte
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(slide, "Introduction & Contexte", "Pourquoi un gestionnaire de mots de passe ?")
accent_line(slide)
footer(slide)

# Chiffres-clés — 3 cartes
stats = [
    ("80 %", "des violations de données\nimpliquent des mots de passe\nfaibles ou volés"),
    ("80+",  "comptes en ligne gérés\nen moyenne par utilisateur"),
    ("12 car.", "minimum recommandé\npar l'ANSSI pour un\nmot de passe sécurisé"),
]
for i, (num, desc) in enumerate(stats):
    cx = Inches(0.5 + i * 4.1)
    add_rect(slide, cx, Inches(1.45), Inches(3.7), Inches(2.4), fill=BLUE_DARK)
    add_textbox(slide, num, cx, Inches(1.55), Inches(3.7), Inches(0.9),
                font_size=34, bold=True, color=RGBColor(0xFF, 0xC1, 0x07),
                align=PP_ALIGN.CENTER)
    add_textbox(slide, desc, cx, Inches(2.45), Inches(3.7), Inches(1.3),
                font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

# Problématique
add_rect(slide, Inches(0.4), Inches(4.1), Inches(12.1), Inches(0.6),
         fill=RGBColor(0xE8, 0xF0, 0xFE))
add_textbox(slide,
            "Problématique : Comment stocker et générer des mots de passe de manière sécurisée, simple et engageante ?",
            Inches(0.5), Inches(4.12), Inches(12), Inches(0.55),
            font_size=14, color=BLUE_DARK, bold=True)

add_bullet_box(slide, [
    "Solution : Application web fullstack conforme ANSSI avec chiffrement AES-256-GCM",
    "Interface moderne inspirée de Proton Pass — accessible sans expertise en cybersécurité",
    "Système de gamification pour encourager l'adoption de bonnes pratiques",
], Inches(0.5), Inches(4.85), Inches(12), Inches(2.0),
   font_size=15, color=BLUE_DARK)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Objectifs
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(slide, "Objectifs du Projet")
accent_line(slide)
footer(slide)

# Colonne gauche — Fonctionnels
add_rect(slide, Inches(0.4), Inches(1.4), Inches(5.9), Inches(5.5), fill=LIGHT_BG)
add_textbox(slide, "⚙  Objectifs fonctionnels",
            Inches(0.5), Inches(1.5), Inches(5.7), Inches(0.5),
            font_size=16, bold=True, color=BLUE_DARK)
func_items = [
    "Inscription / connexion sécurisée",
    "Coffre-fort de mots de passe chiffré",
    "Générateur aléatoire & mémorisable",
    "Catégorisation & recherche",
    "Export CSV / JSON",
    "Historique des mots de passe générés",
]
add_bullet_box(slide, func_items,
               Inches(0.5), Inches(2.05), Inches(5.7), Inches(4.5),
               font_size=14, color=BLUE_DARK)

# Colonne droite — Sécurité
add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.9), Inches(5.5), fill=RGBColor(0xE8, 0xF4, 0xE8))
add_textbox(slide, "🔒  Objectifs de sécurité",
            Inches(6.9), Inches(1.5), Inches(5.7), Inches(0.5),
            font_size=16, bold=True, color=GREEN)
sec_items = [
    "AES-256-GCM pour les mots de passe",
    "bcrypt 12 rounds (mot de passe maître)",
    "JWT — expiration 15 min",
    "Rate limiting : 5 tentatives / 15 min",
    "Politique ANSSI : 12 car. min.",
    "En-têtes sécurité (Helmet.js)",
]
add_bullet_box(slide, sec_items,
               Inches(6.9), Inches(2.05), Inches(5.7), Inches(4.5),
               font_size=14, color=GREEN)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Stack Technologique
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(slide, "Stack Technologique", "Comparaison des approches & choix retenus")
accent_line(slide)
footer(slide)

# Tableau comparatif
cols   = ["Approche", "Exemples", "Pour notre projet ?"]
rows   = [
    ["Site builders (Wix…)", "Wix, Weebly, Jimdo", "✗  Pas de contrôle sécurité"],
    ["CMS",                  "WordPress, Drupal",  "⚠  Surface d'attaque élevée"],
    ["HTML/CSS/JS + Backend","Node.js + Express",  "✓  Contrôle total — CHOIX RETENU"],
]
col_w = [Inches(3.5), Inches(3.0), Inches(5.2)]
col_x = [Inches(0.5), Inches(4.1), Inches(7.2)]

# En-têtes tableau
for j, (col, cx, cw) in enumerate(zip(cols, col_x, col_w)):
    add_rect(slide, cx, Inches(1.45), cw, Inches(0.5), fill=BLUE_DARK)
    add_textbox(slide, col, cx + Pt(4), Inches(1.47), cw, Inches(0.48),
                font_size=13, bold=True, color=WHITE)

row_colors = [RGBColor(0xFF, 0xEE, 0xEE), RGBColor(0xFF, 0xF8, 0xE1),
              RGBColor(0xE8, 0xF5, 0xE9)]
for i, row in enumerate(rows):
    ry = Inches(1.95) + i * Inches(0.65)
    for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_rect(slide, cx, ry, cw, Inches(0.62), fill=row_colors[i],
                 line=RGBColor(0xCC, 0xCC, 0xCC))
        add_textbox(slide, cell, cx + Pt(4), ry + Pt(4), cw, Inches(0.58),
                    font_size=12,
                    color=GREEN if "✓" in cell else (ORANGE if "⚠" in cell else BLUE_DARK))

# Technologies retenues
add_rect(slide, Inches(0.4), Inches(4.0), Inches(12.1), Inches(2.8), fill=LIGHT_BG)
add_textbox(slide, "Technologies retenues",
            Inches(0.5), Inches(4.08), Inches(12), Inches(0.45),
            font_size=15, bold=True, color=BLUE_DARK)

techs = [
    ("Backend", "Node.js 18+  ·  Express 4  ·  Mongoose 8"),
    ("Sécurité", "AES-256-GCM  ·  bcrypt  ·  JWT  ·  Helmet  ·  express-rate-limit"),
    ("Base de données", "MongoDB — 5 collections"),
    ("Frontend", "HTML5  ·  CSS3  ·  JavaScript Vanilla (SPA)"),
    ("Dev / Qualité", "Nodemon  ·  Jest  ·  Supertest  ·  ESLint  ·  Prettier"),
]
for i, (label, val) in enumerate(techs):
    y = Inches(4.58) + i * Inches(0.42)
    add_textbox(slide, f"  {label} :", Inches(0.5), y, Inches(2.2), Inches(0.4),
                font_size=12, bold=True, color=BLUE_DARK)
    add_textbox(slide, val, Inches(2.8), y, Inches(9.5), Inches(0.4),
                font_size=12, color=GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Architecture
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(slide, "Architecture MVC", "Node.js / Express / MongoDB")
accent_line(slide)
footer(slide)

# Diagramme architecture en blocs
blocks = [
    (Inches(0.3),  Inches(1.5),  Inches(3.0),  "🌐  FRONTEND\n\nHTML5 · CSS3 · JS\nSPA publique\n+ Dashboard"),
    (Inches(4.8),  Inches(1.5),  Inches(3.0),  "⚙  API REST\n\nExpress.js\nRoutes · Controllers\nMiddleware · Services"),
    (Inches(9.5),  Inches(1.5),  Inches(3.0),  "🗄  MONGODB\n\nusers\npasswords\ncategories\ngamifications\npasswordhistories"),
]
block_colors = [BLUE, BLUE_DARK, RGBColor(0x1B, 0x5E, 0x20)]
for (bx, by, bw, btxt), bcol in zip(blocks, block_colors):
    add_rect(slide, bx, by, bw, Inches(4.5), fill=bcol)
    add_textbox(slide, btxt, bx + Inches(0.1), by + Inches(0.2),
                bw - Inches(0.2), Inches(4.0),
                font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

# Flèches textuelles
arrow_y = Inches(3.5)
add_textbox(slide, "──  HTTP / JSON  ──►",
            Inches(3.35), arrow_y, Inches(1.4), Inches(0.4),
            font_size=11, color=GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, "──  Mongoose ODM  ──►",
            Inches(7.85), arrow_y, Inches(1.6), Inches(0.4),
            font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Structure fichiers
add_textbox(slide, "Structure du projet :",
            Inches(0.3), Inches(6.1), Inches(6), Inches(0.4),
            font_size=13, bold=True, color=BLUE_DARK)
add_textbox(slide,
            "src/controllers/  ·  src/models/  ·  src/middleware/  ·  src/routes/  ·  src/services/  ·  src/utils/  ·  public/",
            Inches(0.3), Inches(6.55), Inches(12.5), Inches(0.45),
            font_size=11, color=GRAY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Base de données
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(slide, "Base de Données", "MongoDB — 5 collections avec Mongoose ODM")
accent_line(slide)
footer(slide)

collections = [
    ("users", ["_id", "email (unique)", "masterPassword (bcrypt)", "name", "createdAt"]),
    ("passwords", ["_id", "userId (ref)", "title", "username", "encryptedPassword {encrypted, iv, salt, authTag}", "url", "category", "notes"]),
    ("categories", ["_id", "userId (ref)", "name", "icon", "color", "isDefault"]),
    ("gamifications", ["_id", "userId (unique)", "points", "level", "badges[]", "achievements{}", "statistics{}"]),
    ("passwordhistories", ["_id", "userId (ref)", "generatedPassword", "length", "options{}", "strength", "entropy", "used"]),
]
col_positions = [
    (Inches(0.3),  Inches(1.4)),
    (Inches(0.3),  Inches(3.5)),
    (Inches(4.7),  Inches(1.4)),
    (Inches(4.7),  Inches(3.5)),
    (Inches(9.0),  Inches(1.4)),
]
col_width = Inches(4.2)
for (cx, cy), (cname, fields) in zip(col_positions, collections):
    add_rect(slide, cx, cy, col_width, Inches(0.42), fill=BLUE_DARK)
    add_textbox(slide, f"  📄  {cname}", cx + Pt(4), cy + Pt(4),
                col_width, Inches(0.38), font_size=13, bold=True, color=WHITE)
    fh = Inches(0.32) * len(fields)
    add_rect(slide, cx, cy + Inches(0.42), col_width, fh,
             fill=LIGHT_BG, line=BLUE)
    for fi, field in enumerate(fields):
        fy = cy + Inches(0.45) + fi * Inches(0.32)
        add_textbox(slide, f"  · {field}", cx, fy,
                    col_width, Inches(0.3), font_size=10, color=BLUE_DARK)

# Note chiffrement
add_rect(slide, Inches(0.3), Inches(6.55), Inches(12.5), Inches(0.45),
         fill=RGBColor(0xFF, 0xF3, 0xCD))
add_textbox(slide,
            "🔐  Les mots de passe sont stockés chiffrés (AES-256-GCM). Seul le serveur peut les déchiffrer avec la clé de chiffrement.",
            Inches(0.4), Inches(6.58), Inches(12.2), Inches(0.4),
            font_size=11, color=RGBColor(0x6D, 0x4C, 0x00), bold=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Fonctionnalités
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(slide, "Fonctionnalités Clés", "Interface utilisateur & API REST")
accent_line(slide)
footer(slide)

features = [
    ("🔑", "Coffre-fort", "Stockage chiffré\ndes identifiants\n(AES-256-GCM)"),
    ("🔍", "Recherche", "Filtrage temps réel\npar titre, URL,\ncatégorie"),
    ("⚡", "Générateur", "Mode aléatoire\n& mémorisable\navec entropie"),
    ("📁", "Catégories", "8 catégories\npré-définies\n+ personnalisées"),
    ("📤", "Export", "CSV & JSON\navec vérification\ndu mot de passe"),
    ("🏆", "Gamification", "Points, niveaux\net badges de\nsécurité"),
]
for i, (icon, title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    cx = Inches(0.4 + col * 4.2)
    cy = Inches(1.55 + row * 2.6)
    colors_bg = [BLUE_DARK, BLUE, RGBColor(0x1B, 0x5E, 0x20),
                 RGBColor(0x6A, 0x1B, 0x9A), ORANGE, RGBColor(0x0D, 0x47, 0xA1)]
    add_rect(slide, cx, cy, Inches(3.8), Inches(2.3), fill=colors_bg[i])
    add_textbox(slide, icon,
                cx, cy + Inches(0.1), Inches(3.8), Inches(0.55),
                font_size=26, align=PP_ALIGN.CENTER, color=WHITE)
    add_textbox(slide, title,
                cx, cy + Inches(0.65), Inches(3.8), Inches(0.4),
                font_size=15, bold=True, align=PP_ALIGN.CENTER, color=WHITE)
    add_textbox(slide, desc,
                cx, cy + Inches(1.05), Inches(3.8), Inches(1.1),
                font_size=12, align=PP_ALIGN.CENTER, color=RGBColor(0xBB, 0xCC, 0xFF))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Sécurité ANSSI
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(slide, "Sécurité — Conformité ANSSI", "Mécanismes de protection implémentés")
accent_line(slide)
footer(slide)

# Grille 2x3 de mesures de sécurité
sec_measures = [
    ("🔐  AES-256-GCM",
     "Chiffrement authentifié\nPBKDF2 (100k itérations)\nIV aléatoire par entrée\nAuth tag 16 octets"),
    ("🛡  bcrypt — 12 rounds",
     "Hachage du mot de passe maître\n~250 ms/tentative\nProtection brute-force\nSalt automatique"),
    ("🎫  JWT Bearer",
     "Token expiration : 15 min\nPas de sessions côté serveur\nVérification sur chaque route\nHS256 signing"),
    ("⏱  Rate Limiting",
     "5 tentatives / 15 min\nPar adresse IP\nSur l'endpoint /login\nexpress-rate-limit"),
    ("📋  Politique ANSSI",
     "12 caractères minimum\nMajuscule + Minuscule\nChiffre + Symbole\nValidation express-validator"),
    ("🪖  Helmet.js",
     "Content-Security-Policy\nX-Frame-Options: DENY\nX-Content-Type-Options\nHSTS (HTTPS only)"),
]
for i, (stitle, sdesc) in enumerate(sec_measures):
    col = i % 3
    row = i // 3
    cx = Inches(0.3 + col * 4.2)
    cy = Inches(1.45 + row * 2.7)
    add_rect(slide, cx, cy, Inches(4.0), Inches(2.5),
             fill=BLUE_DARK if row == 0 else RGBColor(0x1B, 0x5E, 0x20))
    add_textbox(slide, stitle, cx + Inches(0.1), cy + Inches(0.1),
                Inches(3.8), Inches(0.45), font_size=14, bold=True, color=WHITE)
    add_textbox(slide, sdesc, cx + Inches(0.15), cy + Inches(0.6),
                Inches(3.7), Inches(1.8), font_size=11, color=RGBColor(0xCC, 0xEE, 0xCC))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Générateur de mots de passe
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(slide, "Générateur de Mots de Passe", "Modes aléatoire & mémorisable avec calcul d'entropie")
accent_line(slide)
footer(slide)

# Mode aléatoire
add_rect(slide, Inches(0.3), Inches(1.45), Inches(5.8), Inches(5.3), fill=LIGHT_BG)
add_textbox(slide, "⚡  Mode aléatoire",
            Inches(0.4), Inches(1.55), Inches(5.6), Inches(0.5),
            font_size=16, bold=True, color=BLUE_DARK)
rand_items = [
    "Longueur : 8 à 64 caractères",
    "Sets : minuscules, majuscules,\n     chiffres, symboles",
    "Génération cryptographique\n     (crypto.randomBytes)",
    "Entropie = longueur × log₂(charset)",
    "Exemple (16 car., tous sets) :\n     Entropie ≈ 105 bits",
    "→ Crack time : plusieurs milliards d'années",
]
add_bullet_box(slide, rand_items,
               Inches(0.4), Inches(2.1), Inches(5.5), Inches(4.4),
               font_size=13, color=BLUE_DARK)

# Mode mémorisable
add_rect(slide, Inches(6.6), Inches(1.45), Inches(6.2), Inches(5.3),
         fill=RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, "🧠  Mode mémorisable",
            Inches(6.7), Inches(1.55), Inches(6.0), Inches(0.5),
            font_size=16, bold=True, color=GREEN)
memo_items = [
    "Liste de mots anglais (nature,\n     animaux, technologie, mythologie…)",
    "3 à 8 mots par mot de passe",
    "Séparateurs variés : - . , _ espace",
    "Option capitalisation & chiffres",
    "Exemple :\n     Tiger-Frost-Cipher-42",
    "Plus mémorisable, entropie correcte",
]
add_bullet_box(slide, memo_items,
               Inches(6.7), Inches(2.1), Inches(6.0), Inches(4.4),
               font_size=13, color=GREEN)

# Niveaux de force
add_rect(slide, Inches(0.3), Inches(6.85), Inches(12.5), Inches(0.35),
         fill=BLUE_DARK)
levels = ["< 28 bits : Très faible", "28-36 : Faible", "36-60 : Moyen",
          "60-80 : Fort", "> 80 bits : Très fort ✓"]
lvl_colors = [RGBColor(0xFF, 0x55, 0x55), ORANGE,
              RGBColor(0xFF, 0xC1, 0x07), RGBColor(0x66, 0xBB, 0x6A), WHITE]
for i, (lvl, lc) in enumerate(zip(levels, lvl_colors)):
    add_textbox(slide, lvl, Inches(0.5 + i * 2.4), Inches(6.87),
                Inches(2.3), Inches(0.3), font_size=10, color=lc,
                align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Gamification
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(slide, "Système de Gamification", "Encourager les bonnes pratiques de sécurité")
accent_line(slide)
footer(slide)

add_textbox(slide,
            "La gamification rend la cybersécurité engageante : chaque bonne action est récompensée.",
            Inches(0.5), Inches(1.45), Inches(12), Inches(0.5),
            font_size=14, color=GRAY, italic=True)

# Mécaniques
add_rect(slide, Inches(0.3), Inches(2.05), Inches(5.8), Inches(2.5),
         fill=BLUE_DARK)
add_textbox(slide, "📊  Points & Niveaux",
            Inches(0.4), Inches(2.15), Inches(5.6), Inches(0.45),
            font_size=15, bold=True, color=WHITE)
add_bullet_box(slide, [
    "Points attribués à chaque action",
    "Niveau = floor(points / 100) + 1",
    "Statistiques : total générés, sauvegardés,\n     force moyenne, catégorie favorite",
], Inches(0.4), Inches(2.65), Inches(5.6), Inches(1.8),
   font_size=12, color=WHITE, bullet="→  ")

# Badges
badges_data = [
    ("🔑", "Premier Pas", "1er MDP créé"),
    ("🛡️", "Gardien",     "10 MDP"),
    ("⭐", "Expert",       "50 MDP"),
    ("🔒", "Pro Sécu",     "20 forts"),
    ("📁", "Organisateur", "5 catégories"),
    ("📅", "Fidèle",       "7 jours consécutifs"),
]
add_rect(slide, Inches(6.5), Inches(2.05), Inches(6.2), Inches(2.5),
         fill=RGBColor(0x1B, 0x5E, 0x20))
add_textbox(slide, "🏆  Badges disponibles",
            Inches(6.6), Inches(2.15), Inches(6.0), Inches(0.45),
            font_size=15, bold=True, color=WHITE)
for i, (bicon, bname, bcond) in enumerate(badges_data):
    bx = Inches(6.55 + (i % 3) * 2.0)
    by = Inches(2.68 + (i // 3) * 0.85)
    add_textbox(slide, f"{bicon}  {bname}\n     {bcond}",
                bx, by, Inches(1.9), Inches(0.75),
                font_size=10, color=WHITE)

# Progression
add_rect(slide, Inches(0.3), Inches(4.7), Inches(12.5), Inches(2.0),
         fill=LIGHT_BG)
add_textbox(slide, "Exemple de progression utilisateur :",
            Inches(0.4), Inches(4.78), Inches(12), Inches(0.4),
            font_size=13, bold=True, color=BLUE_DARK)
steps = [
    "Inscription\n0 pts",
    "1er MDP sauvegardé\n+10 pts 🔑",
    "5 MDP forts\n+50 pts",
    "10 MDP créés\n+100 pts 🛡️",
    "Niveau 2\n200 pts ⭐",
]
step_colors = [BLUE, RGBColor(0x0D, 0x47, 0xA1), RGBColor(0x1B, 0x5E, 0x20),
               RGBColor(0x6A, 0x1B, 0x9A), ORANGE]
for i, (step, sc) in enumerate(zip(steps, step_colors)):
    sx = Inches(0.5 + i * 2.42)
    add_rect(slide, sx, Inches(5.25), Inches(2.1), Inches(1.25), fill=sc)
    add_textbox(slide, step, sx, Inches(5.35), Inches(2.1), Inches(1.1),
                font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        add_textbox(slide, "►", Inches(2.5 + i * 2.42), Inches(5.7),
                    Inches(0.2), Inches(0.4), font_size=14, color=GRAY,
                    align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Tests & Qualité
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE)
header_bar(slide, "Tests & Qualité du Code", "Jest · Supertest · ESLint · Prettier")
accent_line(slide)
footer(slide)

# Couverture
add_rect(slide, Inches(0.3), Inches(1.45), Inches(6.0), Inches(3.0),
         fill=LIGHT_BG)
add_textbox(slide, "📊  Seuils de couverture (Jest)",
            Inches(0.4), Inches(1.55), Inches(5.8), Inches(0.45),
            font_size=14, bold=True, color=BLUE_DARK)
coverage = [
    ("Branches",   "80 %"),
    ("Fonctions",  "80 %"),
    ("Lignes",     "85 %"),
    ("Statements", "85 %"),
]
for i, (metric, threshold) in enumerate(coverage):
    my = Inches(2.1) + i * Inches(0.55)
    add_textbox(slide, metric, Inches(0.5), my, Inches(2.5), Inches(0.45),
                font_size=13, color=BLUE_DARK, bold=True)
    # Barre de progression
    add_rect(slide, Inches(3.0), my + Pt(8), Inches(3.0), Inches(0.28),
             fill=RGBColor(0xDD, 0xDD, 0xDD))
    bar_w = Inches(3.0) * (int(threshold.replace(" %", "")) / 100)
    add_rect(slide, Inches(3.0), my + Pt(8), bar_w, Inches(0.28), fill=GREEN)
    add_textbox(slide, threshold, Inches(6.1), my, Inches(0.7), Inches(0.4),
                font_size=12, color=GREEN, bold=True)

# Outils
add_rect(slide, Inches(6.8), Inches(1.45), Inches(5.8), Inches(3.0),
         fill=RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, "🛠  Outils de qualité",
            Inches(6.9), Inches(1.55), Inches(5.6), Inches(0.45),
            font_size=14, bold=True, color=GREEN)
tools_items = [
    "Jest — Tests unitaires et intégration",
    "Supertest — Tests API HTTP",
    "ESLint — Analyse statique du code",
    "Prettier — Formatage automatique",
    "Nodemon — Hot reload en développement",
]
add_bullet_box(slide, tools_items,
               Inches(6.9), Inches(2.1), Inches(5.6), Inches(2.2),
               font_size=13, color=GREEN)

# Git workflow
add_rect(slide, Inches(0.3), Inches(4.6), Inches(12.5), Inches(2.2),
         fill=BLUE_DARK)
add_textbox(slide, "🌿  Git Workflow — Branches",
            Inches(0.4), Inches(4.7), Inches(12), Inches(0.45),
            font_size=14, bold=True, color=WHITE)
branches = [
    ("main", "Production stable"),
    ("dev", "Intégration continue"),
    ("feature/oumar-toure", "Full-Stack Dev"),
    ("feature/rodney", "Security Analyst"),
    ("feature/fresnel", "Backend & Tests"),
]
for i, (bname, brole) in enumerate(branches):
    bx = Inches(0.4 + i * 2.5)
    add_rect(slide, bx, Inches(5.25), Inches(2.2), Inches(1.2),
             fill=BLUE if "main" in bname else (RGBColor(0x2E, 0x7D, 0x32) if "dev" in bname else RGBColor(0x4A, 0x14, 0x8C)))
    add_textbox(slide, f"  {bname}", bx + Pt(4), Inches(5.32),
                Inches(2.1), Inches(0.4), font_size=11, bold=True, color=WHITE)
    add_textbox(slide, brole, bx + Pt(4), Inches(5.75),
                Inches(2.1), Inches(0.55), font_size=10, color=RGBColor(0xCC, 0xBB, 0xFF))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Conclusion & Perspectives
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BLUE_DARK)
add_rect(slide, 0, Inches(4.2), SLIDE_W, Inches(3.3), fill=BLUE)
footer(slide)

add_textbox(slide, "Conclusion & Perspectives",
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
            font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bilan
add_rect(slide, Inches(0.4), Inches(1.2), Inches(5.8), Inches(2.7),
         fill=RGBColor(0x0D, 0x2B, 0x55))
add_textbox(slide, "✅  Bilan du projet",
            Inches(0.5), Inches(1.3), Inches(5.6), Inches(0.45),
            font_size=15, bold=True, color=RGBColor(0xFF, 0xC1, 0x07))
bilan_items = [
    "Architecture MVC sécurisée et modulaire",
    "Chiffrement AES-256-GCM conforme ANSSI",
    "Générateur double-mode avec entropie",
    "Gamification pour l'adoption sécurité",
    "Export et portabilité des données",
]
add_bullet_box(slide, bilan_items,
               Inches(0.5), Inches(1.85), Inches(5.6), Inches(1.9),
               font_size=12, color=WHITE, bullet="✓  ")

# Perspectives
add_rect(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.7),
         fill=RGBColor(0x0D, 0x2B, 0x55))
add_textbox(slide, "🚀  Perspectives",
            Inches(6.9), Inches(1.3), Inches(5.6), Inches(0.45),
            font_size=15, bold=True, color=RGBColor(0xFF, 0xC1, 0x07))
persp_items = [
    "Authentification 2FA (TOTP/SMS)",
    "Extension navigateur (auto-fill)",
    "Architecture Zero-Knowledge",
    "Application mobile (React Native)",
    "Partage sécurisé de mots de passe",
]
add_bullet_box(slide, persp_items,
               Inches(6.9), Inches(1.85), Inches(5.6), Inches(1.9),
               font_size=12, color=WHITE, bullet="→  ")

# Message final
add_textbox(slide,
            "TAMYCS Shield démontre qu'il est possible de rendre la sécurité\naccessible, engageante et techniquement robuste.",
            Inches(0.5), Inches(4.35), Inches(12), Inches(0.9),
            font_size=16, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

add_textbox(slide, "Merci pour votre attention  🛡",
            Inches(0.5), Inches(5.5), Inches(12), Inches(0.7),
            font_size=28, bold=True, color=RGBColor(0xFF, 0xC1, 0x07),
            align=PP_ALIGN.CENTER)

add_textbox(slide,
            "Oumar Touré  ·  Rodney  ·  Fresnel\nESAIP La Salle — Bachelor 3 Cybersécurité — 2025/2026",
            Inches(0.5), Inches(6.2), Inches(12), Inches(0.7),
            font_size=13, color=RGBColor(0xAA, 0xBB, 0xDD),
            align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# Sauvegarde
# ─────────────────────────────────────────────────────────────────────────────
output = os.path.join(os.path.dirname(__file__), "TAMYCS_Shield_Presentation.pptx")
prs.save(output)
print("PowerPoint cree : " + output)
print("Nombre de slides : " + str(len(prs.slides)))
